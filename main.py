from pptx import Presentation
from pptx.util import Inches
from wand.image import Image
from PIL import Image as pil_img
import os


def logo_search(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith('.png'):
            img = filename
            image = pil_img.open(folder_path + '/' + filename)
            new_image = image.resize((image.width // 2, image.height // 2))  # resizing logo
            new_image.save('./altered_images/resized_img.png')  # saving resized logo
            return img


def watermark(folder_path):
    prs = Presentation()  # Presentation Object

    # Searching for logo
    img = logo_search(folder_path)

    i = 1
    for filename in os.listdir(folder_path):
        if filename != img:
            # importing the images that has to be watermarked
            with Image(filename=folder_path + '/' + filename) as image:

                # Import the watermark image
                with Image(filename='./altered_images/resized_img.png') as water:

                    # Clone the image in order to process
                    with image.clone() as watermark:

                        # Invoke watermark function with watermark image
                        # left as 100 and top as 50
                        watermark.watermark(water, 0, 100, 50)
                        watermark.save(filename='./altered_images/' + filename)  # saving the image

                        title_slide_layout = prs.slide_layouts[1]  # creating the layout 1 in presentation
                        slide = prs.slides.add_slide(title_slide_layout)  # adding slide to the presentation
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        left = Inches(1)
                        top = Inches(2.5)
                        title.text = "Sample Title " + str(i)  # adding title to ppt
                        subtitle.text = "Sample Subtitle " + str(i)  # adding subtitle to ppt
                        i += 1
                        height = Inches(3.5)  # setting the height of the image in ppt

                        slide.shapes.add_picture('./altered_images/' + filename, left, top, height=height)  # adding
                        # picture to ppt
                        prs.save('test.pptx')  # saving the ppt


watermark('./PPT - Logo - Images - Assignment')


